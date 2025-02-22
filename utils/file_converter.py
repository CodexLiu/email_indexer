import asyncio
import aiohttp
from pathlib import Path
from typing import Optional, List, Dict
from docx import Document
import pandas as pd
from pdf2image import convert_from_path
import io
from PIL import Image
import base64
from concurrent.futures import ThreadPoolExecutor
from utils.openrouter_utils import make_image_api_call
from utils.prompts import pdf_transcription_prompt
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


async def convert_file_to_text(file_path: str) -> Optional[str]:
    """Convert any supported file to text."""
    file_path = Path(file_path)
    extension = file_path.suffix.lower()

    try:
        if extension == '.pdf':
            result = await process_single_pdf(str(file_path))
            if result:
                return result.get(file_path.stem, "")
            return None

        elif extension == '.docx':
            return await process_docx_file(str(file_path))

        elif extension in ['.xlsx', '.xls', '.xlsm']:
            text_content = read_excel_file(str(file_path))
            if not text_content:
                return None
            return text_content

        elif extension == '.pptx':
            result = await process_pptx_file(str(file_path))
            if result:
                return result.get(file_path.stem, "")
            return None

        else:
            return None

    except Exception as e:
        print(f"Error converting file {file_path}: {str(e)}")
        return None


async def process_docx_file(file_path: str) -> Optional[str]:
    """Extract text content from a DOCX file."""
    try:
        # Run in thread pool since this is CPU-bound
        loop = asyncio.get_event_loop()
        with ThreadPoolExecutor() as pool:
            doc = await loop.run_in_executor(pool, Document, file_path)
            paragraphs = [paragraph.text for paragraph in doc.paragraphs]
            return "\n".join(paragraphs)
    except Exception as e:
        print(f"Error processing DOCX file: {str(e)}")
        return None


def read_excel_file(file_path: str) -> str:
    """Read Excel file and convert all sheets to text format."""
    try:
        excel = pd.ExcelFile(file_path)
        all_text = []

        for sheet_name in excel.sheet_names:
            try:
                df = pd.read_excel(
                    file_path,
                    sheet_name=sheet_name,
                    engine='openpyxl'
                )

                sheet_text = f"\n{'='*20} Sheet: {sheet_name} {'='*20}\n"

                if df.empty:
                    sheet_text += "\n[Empty Sheet]\n"
                else:
                    pd.set_option('display.max_colwidth', None)
                    pd.set_option('display.max_columns', None)
                    pd.set_option('display.width', None)

                    sheet_content = df.fillna('').to_string(
                        index=False,
                        max_rows=None,
                        max_cols=None
                    )
                    sheet_text += sheet_content

                all_text.append(sheet_text)

            except Exception as e:
                all_text.append(
                    f"\nError reading sheet '{sheet_name}': {str(e)}\n")

        return "\n\n" + "\n\n".join(all_text) + "\n"

    except Exception as e:
        print(f"Error reading Excel file: {str(e)}")
        return ""


def convert_pdf_to_images(pdf_path):
    """Convert PDF to images in a thread-safe manner."""
    return convert_from_path(pdf_path)


async def process_single_page(session: aiohttp.ClientSession,
                              image: Image.Image,
                              page_num: int,
                              max_retries=5,
                              base_delay=1) -> tuple[int, str, bool]:
    """Process a single page with retry logic."""
    print(f"Processing page {page_num}...")

    # Convert image to base64
    img_byte_arr = io.BytesIO()
    image.save(img_byte_arr, format='JPEG')
    img_byte_arr = img_byte_arr.getvalue()
    img_base64 = base64.b64encode(img_byte_arr).decode('utf-8')

    # Get API call parameters
    api_params = make_image_api_call(img_base64, pdf_transcription_prompt)

    # Modify the API call to use the existing session
    for attempt in range(max_retries):
        try:
            async with session.post(api_params['url'],
                                    json=api_params['json'],
                                    headers=api_params['headers']) as response:
                if response.status == 200:
                    result = await response.json()
                    if 'choices' in result and result['choices']:
                        transcribed_text = result['choices'][0]['message']['content']
                        print(f"Successfully processed page {page_num}")
                        return page_num, transcribed_text, True
                elif response.status == 429:  # Rate limit exceeded
                    delay = base_delay * (2 ** attempt)
                    print(
                        f"Rate limit exceeded for page {page_num}. Retrying in {delay} seconds...")
                    await asyncio.sleep(delay)
                    continue

                # If we get here, the request failed
                delay = base_delay * (2 ** attempt)
                print(
                    f"Attempt {attempt + 1} for page {page_num} failed (status {response.status}). Retrying in {delay} seconds...")
                await asyncio.sleep(delay)

        except Exception as e:
            if attempt == max_retries - 1:
                print(
                    f"All retry attempts failed for page {page_num}. Final error: {str(e)}")
                return page_num, f"[Failed to process page {page_num} after {max_retries} attempts]", False

            delay = base_delay * (2 ** attempt)
            print(
                f"Error processing page {page_num}: {str(e)}. Retrying in {delay} seconds...")
            await asyncio.sleep(delay)

    return page_num, f"[Failed to process page {page_num} after {max_retries} attempts]", False


async def process_single_pdf(pdf_path: str, max_concurrent=100) -> dict[str, str]:
    """Process a single PDF file with limited concurrency."""
    try:
        base_filename = Path(pdf_path).stem

        # Convert PDF to images in a thread pool since it's CPU-bound
        with ThreadPoolExecutor() as pool:
            images = await asyncio.get_event_loop().run_in_executor(
                pool, convert_pdf_to_images, pdf_path
            )

        async with aiohttp.ClientSession() as session:
            # Process pages in batches to limit concurrency
            results = []
            for i in range(0, len(images), max_concurrent):
                batch = images[i:i + max_concurrent]

                # Create tasks for current batch
                tasks = []
                for j, image in enumerate(batch, i + 1):
                    task = asyncio.create_task(
                        process_single_page(session, image, j))
                    tasks.append(task)

                # Process current batch concurrently
                batch_results = await asyncio.gather(*tasks, return_exceptions=True)

                # Handle results
                for result in batch_results:
                    if isinstance(result, Exception):
                        print(f"Error in batch processing: {str(result)}")
                        results.append(
                            (len(results) + 1, "[Page processing failed]", False))
                    else:
                        results.append(result)

            # Sort and combine results
            complete_text = ""
            results.sort(key=lambda x: x[0])  # Sort by page number

            for i, (page_num, text, success) in enumerate(results, 1):
                complete_text += f"\n{'='*50}\nPage {i}\n{'='*50}\n\n"
                complete_text += text if success else f"[Failed to process page {i}]"
                complete_text += '\n\n'

            if not any(success for _, _, success in results):
                print(f"Warning: All pages failed to process in {pdf_path}")

            return {base_filename: complete_text}

    except Exception as e:
        print(f"Error processing PDF {pdf_path}: {str(e)}")
        return {base_filename: ""}


async def process_pptx_file(file_path: str) -> Dict[str, str]:
    """Process PPTX file by converting slides to images and transcribing."""
    try:
        # Run PowerPoint operations in thread pool since they're CPU-bound
        loop = asyncio.get_event_loop()
        with ThreadPoolExecutor() as pool:
            presentation = await loop.run_in_executor(
                pool, Presentation, file_path)
            base_filename = Path(file_path).stem
            all_text = []

            # Process each slide
            async with aiohttp.ClientSession() as session:
                for slide_number, slide in enumerate(presentation.slides, 1):
                    # Convert slide to image
                    image = await loop.run_in_executor(
                        pool, convert_slide_to_image, slide)

                    if image:
                        # Convert image to base64
                        img_byte_arr = io.BytesIO()
                        image.save(img_byte_arr, format='PNG')
                        img_base64 = base64.b64encode(
                            img_byte_arr.getvalue()).decode('utf-8')

                        # Get API call parameters for transcription
                        api_params = make_image_api_call(
                            img_base64, pdf_transcription_prompt)

                        # Process slide with retries
                        _, slide_text, _ = await process_single_page(
                            session, image, slide_number)
                        all_text.append(
                            f"\n{'='*50}\nSlide {slide_number}\n{'='*50}\n\n{slide_text}")

            return {base_filename: "\n".join(all_text)}

    except Exception as e:
        print(f"Error processing PPTX {file_path}: {str(e)}")
        return {base_filename: ""}


def convert_slide_to_image(slide) -> Optional[Image.Image]:
    """Convert a PowerPoint slide to PIL Image."""
    try:
        # Get slide dimensions
        width = int(slide.slide_width)
        height = int(slide.slide_height)

        # Create blank image
        image = Image.new('RGB', (width, height), 'white')

        # Draw shapes on image
        for shape in slide.shapes:
            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                # Handle images in the slide
                image_stream = io.BytesIO(shape.image.blob)
                img = Image.open(image_stream)
                # Calculate position and paste
                left = shape.left
                top = shape.top
                image.paste(img, (left, top))

        return image
    except Exception as e:
        print(f"Error converting slide to image: {str(e)}")
        return None
